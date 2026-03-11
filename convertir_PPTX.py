from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

def create_presentation():
    prs = Presentation()
    
    # Definición de Colores (Basados en el código Slate/Blue/Red)
    SLATE_BG = RGBColor(248, 250, 252)  # #f8fafc [cite: 10]
    SLATE_TEXT = RGBColor(51, 65, 85)   # #334155 [cite: 10]
    BLUE_BRAND = RGBColor(59, 130, 246) # #3b82f6 [cite: 14]
    RED_ALERT = RGBColor(239, 68, 68)   # #ef4444 [cite: 18, 107]
    EMERALD_SUCCESS = RGBColor(16, 185, 129) # #10b981 

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = SLATE_BG

    # --- DIAPOSITIVA 1: DASHBOARD ---
    slide_layout = prs.slide_layouts[6] # Layout en blanco
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # Título Principal [cite: 25, 29]
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Logística Planta - APT: Interaxia"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = SLATE_TEXT

    # Gráfico de Impacto (Doughnut) [cite: 31, 107]
    chart_data = ChartData()
    chart_data.categories = ['Despachos Exitosos', 'Retenidos (Falta Doc)']
    chart_data.add_series('Estatus', (0.82, 0.18))

    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(4.5), Inches(4)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM

    # Hallazgo Crítico [cite: 35, 36]
    shape = slide.shapes.add_shape(6, Inches(5.5), Inches(2), Inches(3.5), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED_ALERT
    tf = shape.text_frame
    tf.text = "Hallazgo Crítico: Movimientos 'Sin Documentación' paralizan ingresos."
    
    # --- DIAPOSITIVA 2: FLUJO OPERATIVO ---
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    
    # Título [cite: 39]
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_box.text_frame.text = "Línea de Tiempo del Proceso"
    
    # Pasos del Flujo [cite: 66, 73]
    steps = [
        "1. Zona Tránsito", "2. Cargar Camión", "3. VERIFICACIÓN (CRÍTICO)", 
        "4. Envío a APT", "5. Pistoleo", "6. Almacenaje"
    ]
    
    for i, step in enumerate(steps):
        box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2 + (i*0.8)), Inches(8), Inches(0.6))
        box.fill.solid()
        # Resaltar paso 3 en rojo 
        box.fill.fore_color.rgb = RED_ALERT if "3" in step else BLUE_BRAND
        box.text_frame.text = step

    # --- DIAPOSITIVA 3: PLAN DE MEJORA ---
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    
    # Título [cite: 62]
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_box.text_frame.text = "Plan de Mejora Estratégica"
    
    # Mejoras [cite: 74, 81]
    mejoras = [
        "Checklist Digital: Bloqueo automático si falta documentación.",
        "Alertas Automáticas: Visibilidad total en tiempo real.",
        "Estandard Etiquetas: 100% legible para escáneres."
    ]
    
    for i, mejora in enumerate(mejoras):
        box = slide.shapes.add_textbox(Inches(1), Inches(2 + (i*1.2)), Inches(8), Inches(1))
        box.text_frame.word_wrap = True
        p = box.text_frame.paragraphs[0]
        p.text = f"• {mejora}"
        p.font.size = Pt(18)
        p.font.color.rgb = SLATE_TEXT

    prs.save('Presentacion_Interaxia_Logistica.pptx')
    print("Archivo 'Presentacion_Interaxia_Logistica.pptx' generado con éxito.")

if __name__ == "__main__":
    create_presentation()