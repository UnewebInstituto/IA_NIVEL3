from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

def create_interaxia_presentation():
    prs = Presentation()
    
    # Paleta de colores Interaxia extraída del código fuente
    SLATE_BG = RGBColor(248, 250, 252)    # Fondo principal [cite: 9]
    SLATE_TEXT = RGBColor(51, 65, 85)     # Texto [cite: 10]
    BLUE_BRAND = RGBColor(37, 99, 235)    # Primario [cite: 13]
    RED_ALERT = RGBColor(220, 38, 38)     # Alertas [cite: 18]
    NAVY_PANEL = RGBColor(15, 23, 42)     # Fondo de análisis (img1)

    def apply_base_style(slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = SLATE_BG
        
        # Header (Título común en img1 e img2)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Gestión de Traslado a APT"
        p.font.bold = True
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(15, 23, 42)
        
        sub = tf.add_paragraph()
        sub.text = "Digitalización de Procesos y Normativas"
        sub.font.size = Pt(12)
        sub.font.color.rgb = RGBColor(100, 116, 139)

    # --- DIAPOSITIVA 1: LÍNEA DE TIEMPO (Basada en img1) ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_base_style(slide1)
    
    # Datos de los 6 pasos (de img1 y flowData) [cite: 66-73]
    pasos = [
        ("01", "Zona de Tránsito", "Origen en Planta"),
        ("02", "Cargar Camión", "¿Docs listos?"),
        ("03", "Verificación", "¡Sin Documentación!", True), # Estado crítico
        ("04", "Envío a APT", "Traslado físico"),
        ("05", "Sistema", "Pistoleo / Registro"),
        ("06", "Almacenamiento", "Ubicación APT")
    ]
    
    for i, (num, titulo, desc, *is_crit) in enumerate(pasos):
        x_pos = Inches(0.5 + (i * 1.5))
        # Tarjeta blanca
        card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(1.5), Inches(1.4), Inches(2))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RED_ALERT if is_crit else RGBColor(226, 232, 240)
        
        # Contenido de la tarjeta
        tf = card.text_frame
        tf.margin_top = Inches(0.1)
        p1 = tf.paragraphs[0]
        p1.text = num
        p1.alignment = PP_ALIGN.RIGHT
        p1.font.size = Pt(10)
        p1.font.color.rgb = RGBColor(148, 163, 184)
        
        p2 = tf.add_paragraph()
        p2.text = f"\n{titulo}"
        p2.font.bold = True
        p2.font.size = Pt(11)
        p2.font.color.rgb = SLATE_TEXT
        
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(9)
        if is_crit:
            p3.font.color.rgb = RED_ALERT
            p3.font.bold = True

    # Panel de Análisis del Experto (img1 inferior derecha) [cite: 42]
    analisis = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(4.2), Inches(4.3), Inches(2.3))
    analisis.fill.solid()
    analisis.fill.fore_color.rgb = NAVY_PANEL
    tf_a = analisis.text_frame
    tf_a.word_wrap = True
    p_a = tf_a.paragraphs[0]
    p_a.text = "🔑 Análisis del Experto"
    p_a.font.bold = True
    p_a.font.color.rgb = RGBColor(255, 255, 255)
    
    p_a2 = tf_a.add_paragraph()
    p_a2.text = "El punto de decisión de Documentación es el cuello de botella identificado. Sin un flujo de papel o digital previo, el ingreso al sistema en APT se retrasa."
    p_a2.font.size = Pt(10)
    p_a2.font.color.rgb = RGBColor(226, 232, 240)

    # --- DIAPOSITIVA 2: DIAGRAMA DE FLUJO (Basada en img2) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_base_style(slide2)
    
    # Nodos del flujo vertical (de img2)
    nodos = [
        "Inicio (ZONA TRÁNSITO)",
        "Cargar Camión (OPERACIONES)",
        "¿Documentos OK? (CHECKLIST)",
        "Envío a APT (TRANSPORTE)",
        "Almacenamiento (UBICACIÓN)"
    ]
    
    for i, texto in enumerate(nodos):
        y_pos = Inches(1.2 + (i * 1.0))
        shape_type = MSO_SHAPE.DIAMOND if "?" in texto else MSO_SHAPE.ROUNDED_RECTANGLE
        
        node = slide2.shapes.add_shape(shape_type, Inches(3.5), y_pos, Inches(3), Inches(0.7))
        node.fill.solid()
        node.fill.fore_color.rgb = RGBColor(255, 255, 255)
        node.line.color.rgb = BLUE_BRAND
        
        p_n = node.text_frame.paragraphs[0]
        p_n.text = texto
        p_n.font.size = Pt(10)
        p_n.font.color.rgb = SLATE_TEXT
        p_n.alignment = PP_ALIGN.CENTER

    # Footer para ambas diapositivas [cite: 64]
    for s in prs.slides:
        footer = s.shapes.add_textbox(Inches(0), Inches(7.2), Inches(10), Inches(0.3))
        footer.fill.solid()
        footer.fill.fore_color.rgb = RGBColor(30, 41, 59)
        p_f = footer.text_frame.paragraphs[0]
        p_f.text = "Interaxia | Panel de Gestión y Trazabilidad Documental 2024"
        p_f.font.size = Pt(9)
        p_f.font.color.rgb = RGBColor(255, 255, 255)
        p_f.alignment = PP_ALIGN.CENTER

    prs.save('Presentacion_Interaxia_Completa.pptx')
    print("Presentación generada exitosamente.")

if __name__ == "__main__":
    create_interaxia_presentation()